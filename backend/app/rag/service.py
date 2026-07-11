# -*- coding: utf-8 -*-
"""RAG service: document parsing, chunking, and retrieval.
Supports both InMemoryStore (default) and Qdrant vector store (when available).
"""
import re
from pathlib import Path
from typing import Optional

from app.config import settings
from app.logging_config import logger
from app.models.document import Document as DocModel

CHUNK_SIZE = 500
CHUNK_OVERLAP = 100


def split_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    if not text.strip():
        return []
    chunks = []
    start = 0
    text_len = len(text)

    while start < text_len:
        if start < 0:
            start = 0
        end = min(start + chunk_size, text_len)
        if end < text_len:
            last_period = text.rfind(".", start, end)
            if last_period > start + chunk_size // 2:
                end = last_period + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        next_start = end - overlap
        if next_start <= start:
            next_start = end
        start = next_start

    return [c for c in chunks if len(c) > 20]


def parse_text(content: str) -> str:
    """Clean and normalize text content"""
    content = re.sub(r"\s+", " ", content)
    return content.strip()


def parse_pdf(file_path: Path) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(file_path))
        text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return "\n".join(text)
    except ImportError:
        raise ImportError("PyPDF2 is required for PDF parsing")


def parse_docx(file_path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        raise ImportError("python-docx is required for DOCX parsing")


def parse_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    except ImportError:
        raise ImportError("python-pptx is required for PPTX parsing")


def parse_document(file_path: Path, file_type: str) -> str:
    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "pptx": parse_pptx,
        "txt": lambda p: p.read_text(encoding="utf-8", errors="replace"),
        "md": lambda p: p.read_text(encoding="utf-8", errors="replace"),
    }
    parser = parsers.get(file_type)
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")
    return parser(file_path)


# ===== InMemoryStore (fallback) =====

class InMemoryStore:
    """Simple in-memory chunk storage for when Qdrant is unavailable."""

    def __init__(self):
        self.chunks: dict[int, list[dict]] = {}
        self.doc_courses: dict[int, int] = {}

    def add_chunks(self, doc_id: int, chunks: list[str], course_id: int = 0) -> None:
        self.chunks[doc_id] = [
            {"id": i, "doc_id": doc_id, "text": chunk}
            for i, chunk in enumerate(chunks)
        ]
        if course_id:
            self.doc_courses[doc_id] = course_id

    def search(self, query: str, top_k: int = 5) -> list[dict]:
        query_terms = set(query.lower().split())
        results = []
        for doc_id, chunks in self.chunks.items():
            for chunk in chunks:
                chunk_terms = set(chunk["text"].lower().split())
                overlap = len(query_terms & chunk_terms)
                if overlap > 0:
                    results.append((overlap, chunk))
        results.sort(key=lambda x: x[0], reverse=True)
        return [r[1] for r in results[:top_k]]

    def get_course_chunks(self, course_id: int, limit: int = 10) -> list[dict]:
        doc_ids = [d for d, c in self.doc_courses.items() if c == course_id]
        all_chunks = []
        for doc_id in doc_ids:
            all_chunks.extend(self.chunks.get(doc_id, []))
        return all_chunks[:limit]

    def remove_document(self, doc_id: int) -> None:
        self.chunks.pop(doc_id, None)
        self.doc_courses.pop(doc_id, None)


# Global stores
store = InMemoryStore()
_use_qdrant = False


async def init_stores():
    """Try to initialize Qdrant; fall back to InMemoryStore."""
    global _use_qdrant
    try:
        from app.rag import qdrant_store as qs
        from app.rag.embedding import get_embedding_dim

        if await qs.is_available():
            dim = await get_embedding_dim()
            await qs.ensure_collection(dim=dim)
            _use_qdrant = True
            logger.info(f"Qdrant ready (dim={dim})")
        else:
            logger.info("Qdrant not available, using InMemoryStore")
    except Exception as e:
        logger.warning(f"Qdrant init failed (see detail): {e}")
        import traceback
        for line in traceback.format_exc().split(chr(10)):
            if line.strip():
                logger.error(f"  {line}")


async def reprocess_existing_documents():
    """Re-process all ready documents to rebuild chunk stores after restart."""
    try:
        from app.database import async_session
        from app.models.document import Document as DocModel
        from sqlalchemy import select
        
        upload_dir = Path(settings.upload_dir)
        async with async_session() as db:
            result = await db.execute(
                select(DocModel).where(DocModel.status == "ready", DocModel.file_size > 0)
            )
            docs = result.scalars().all()
            
        if not docs:
            logger.info("No existing documents to reprocess")
            return
            
        count = 0
        for doc in docs:
            file_path = upload_dir / doc.filename
            if not file_path.exists():
                continue
            try:
                await process_document(doc, file_path)
                count += 1
            except Exception as e:
                logger.warning(f"Reprocess doc {doc.id} ({doc.original_filename}) failed: {e}")
                
        logger.info(f"Reprocessed {count}/{len(docs)} existing documents")
    except Exception as e:
        logger.warning(f"Document reprocess failed: {e}")


async def process_document(doc: DocModel, file_path: Path) -> list[str]:
    """Process a document: parse, chunk, and store."""
    raw_text = parse_document(file_path, doc.file_type)
    clean_text = parse_text(raw_text)
    chunks = split_text(clean_text)

    # Store in InMemoryStore (always)
    store.add_chunks(doc.id, chunks, course_id=doc.course_id)

    # Also store in Qdrant if available
    if _use_qdrant:
        from app.rag import qdrant_store as qs
        await qs.add_chunks(doc.id, chunks, doc.course_id)

    return chunks


async def retrieve_context(query: str, course_id: int, top_k: int = 5) -> str:
    """Retrieve context for a query using available stores."""
    results = []

    if _use_qdrant:
        from app.rag import qdrant_store as qs
        results = await qs.search(query, top_k=top_k, course_id=course_id)

    if not results:
        results = store.search(query, top_k=top_k)

    if not results:
        results = store.get_course_chunks(course_id, limit=top_k)

    if not results:
        return ""

    parts = []
    for i, chunk in enumerate(results, 1):
        text = chunk.get("text", str(chunk))
        parts.append(f"[{i}] {text}")
    return "\n\n".join(parts)


